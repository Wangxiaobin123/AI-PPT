"""PowerPoint presentation generator using python-pptx."""

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .base import (
    BaseGenerator,
    ContentBlock,
    ContentType,
    GenerationMetadata,
    SlideData,
    StructuredContent,
)


# Default color palette
DEFAULT_COLORS = {
    "primary": "2B579A",       # Dark blue
    "secondary": "217346",     # Green
    "accent": "B7472A",        # Red accent
    "background": "FFFFFF",    # White
    "text": "333333",          # Dark gray
    "subtitle": "666666",      # Medium gray
    "table_header": "2B579A",  # Same as primary
    "table_alt_row": "F2F2F2", # Light gray
}

DEFAULT_FONTS = {
    "title": "Calibri",
    "body": "Calibri",
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a hex color string (without #) to an RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _resolve_colors(metadata: GenerationMetadata) -> dict[str, str]:
    """Merge user-supplied color_scheme with defaults."""
    colors = dict(DEFAULT_COLORS)
    user_scheme = metadata.style.get("color_scheme", {})
    colors.update(user_scheme)
    return colors


def _resolve_fonts(metadata: GenerationMetadata) -> dict[str, str]:
    """Merge user-supplied fonts with defaults."""
    fonts = dict(DEFAULT_FONTS)
    user_fonts = metadata.style.get("fonts", {})
    fonts.update(user_fonts)
    return fonts


def _set_text_style(
    run,
    *,
    font_name: str = "Calibri",
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
):
    """Apply common text styling to a run."""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = _hex_to_rgb(color)


class PPTXGenerator(BaseGenerator):
    """Generates PowerPoint (.pptx) presentations from structured content."""

    async def generate(
        self, content: StructuredContent, metadata: GenerationMetadata
    ) -> bytes:
        """Generate a PPTX file and return it as bytes."""
        prs = Presentation()

        # Slide dimensions: standard widescreen 13.333 x 7.5 inches
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        colors = _resolve_colors(metadata)
        fonts = _resolve_fonts(metadata)

        if not content.slides:
            # If no slides provided, create a single title slide
            self._add_title_slide(prs, content.title, content.subtitle, colors, fonts)
        else:
            for slide_data in content.slides:
                layout = slide_data.layout.lower()
                if layout == "title":
                    self._add_title_slide(
                        prs,
                        slide_data.title or content.title,
                        slide_data.subtitle or content.subtitle,
                        colors,
                        fonts,
                    )
                elif layout == "two_column":
                    self._add_two_column_slide(prs, slide_data, colors, fonts)
                elif layout == "blank":
                    self._add_blank_slide(prs)
                else:
                    # "content" or "image" layout
                    self._add_content_slide(prs, slide_data, colors, fonts)

                # Add speaker notes if present
                if slide_data.notes:
                    slide = prs.slides[-1]
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data.notes

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------ #
    #  Slide builders
    # ------------------------------------------------------------------ #

    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str,
        colors: dict[str, str],
        fonts: dict[str, str],
    ):
        """Add a title slide with centered title and subtitle."""
        slide_layout = prs.slide_layouts[6]  # Blank layout for full control
        slide = prs.slides.add_slide(slide_layout)

        # Background
        self._set_slide_background(slide, colors.get("background", "FFFFFF"))

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Title text box – centered in top half
        title_left = Inches(1)
        title_top = Inches(1.8)
        title_width = slide_w - Inches(2)
        title_height = Inches(2)
        txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        _set_text_style(
            run, font_name=fonts["title"], font_size=44, bold=True, color=colors["primary"]
        )

        # Subtitle text box
        if subtitle:
            sub_top = title_top + title_height + Inches(0.2)
            sub_height = Inches(1.2)
            txBox2 = slide.shapes.add_textbox(title_left, sub_top, title_width, sub_height)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = subtitle
            _set_text_style(
                run2,
                font_name=fonts["body"],
                font_size=24,
                color=colors["subtitle"],
            )

        # Decorative bottom bar
        bar_height = Inches(0.08)
        bar_top = slide_h - Inches(0.8)
        bar_left = Inches(3)
        bar_width = slide_w - Inches(6)
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            bar_left,
            bar_top,
            bar_width,
            bar_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(colors["primary"])
        shape.line.fill.background()

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: SlideData,
        colors: dict[str, str],
        fonts: dict[str, str],
    ):
        """Add a standard content slide with title and content blocks."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors.get("background", "FFFFFF"))

        slide_w = prs.slide_width

        # Title bar at the top
        if slide_data.title:
            title_left = Inches(0.5)
            title_top = Inches(0.3)
            title_width = slide_w - Inches(1)
            title_height = Inches(0.9)
            txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = slide_data.title
            _set_text_style(
                run, font_name=fonts["title"], font_size=36, bold=True, color=colors["primary"]
            )

            # Underline accent
            line_top = title_top + title_height
            line_shape = slide.shapes.add_shape(
                1, Inches(0.5), line_top, Inches(3), Inches(0.04)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = _hex_to_rgb(colors["accent"])
            line_shape.line.fill.background()

        # Content area
        content_top = Inches(1.5)
        content_left = Inches(0.7)
        content_width = slide_w - Inches(1.4)
        current_top = content_top

        for block in slide_data.blocks:
            current_top = self._render_block(
                slide, block, content_left, current_top, content_width, colors, fonts
            )

    def _add_two_column_slide(
        self,
        prs: Presentation,
        slide_data: SlideData,
        colors: dict[str, str],
        fonts: dict[str, str],
    ):
        """Add a two-column layout slide. Blocks split evenly between columns."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors.get("background", "FFFFFF"))

        slide_w = prs.slide_width

        # Title
        if slide_data.title:
            title_left = Inches(0.5)
            title_top = Inches(0.3)
            title_width = slide_w - Inches(1)
            title_height = Inches(0.9)
            txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = slide_data.title
            _set_text_style(
                run, font_name=fonts["title"], font_size=36, bold=True, color=colors["primary"]
            )

        # Split blocks into two columns
        blocks = slide_data.blocks
        mid = (len(blocks) + 1) // 2
        left_blocks = blocks[:mid]
        right_blocks = blocks[mid:]

        col_gap = Inches(0.4)
        col_width = (slide_w - Inches(1.4) - col_gap) / 2
        # Convert to int (EMU) for pptx
        col_width_emu = int(col_width)

        left_x = Inches(0.7)
        right_x = left_x + col_width_emu + int(col_gap)
        content_top = Inches(1.5)

        # Left column
        current_top = content_top
        for block in left_blocks:
            current_top = self._render_block(
                slide, block, left_x, current_top, col_width_emu, colors, fonts
            )

        # Right column
        current_top = content_top
        for block in right_blocks:
            current_top = self._render_block(
                slide, block, right_x, current_top, col_width_emu, colors, fonts
            )

    def _add_blank_slide(self, prs: Presentation):
        """Add an empty blank slide."""
        slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(slide_layout)

    # ------------------------------------------------------------------ #
    #  Block renderers
    # ------------------------------------------------------------------ #

    def _render_block(
        self,
        slide,
        block: ContentBlock,
        left,
        top,
        width,
        colors: dict[str, str],
        fonts: dict[str, str],
    ) -> Any:
        """Render a single content block and return the new y-position."""
        if block.type == ContentType.TABLE:
            return self._render_table(slide, block, left, top, width, colors, fonts)
        elif block.type in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
            return self._render_list(slide, block, left, top, width, colors, fonts)
        elif block.type == ContentType.HEADING:
            return self._render_heading(slide, block, left, top, width, colors, fonts)
        elif block.type == ContentType.CODE:
            return self._render_code(slide, block, left, top, width, colors, fonts)
        else:
            # TEXT, IMAGE placeholders, CHART placeholders
            return self._render_text(slide, block, left, top, width, colors, fonts)

    def _render_text(self, slide, block, left, top, width, colors, fonts) -> int:
        """Render a text block and return the updated y-position."""
        text = block.content
        if not text:
            return top

        est_lines = max(1, len(text) // 90 + 1)
        height = Inches(0.35 * est_lines)
        txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_text_style(run, font_name=fonts["body"], font_size=16, color=colors["text"])
        return int(top) + int(height) + Inches(0.15)

    def _render_heading(self, slide, block, left, top, width, colors, fonts) -> int:
        """Render a heading block."""
        text = block.content
        if not text:
            return top

        # Map heading level to font size
        sizes = {1: 32, 2: 28, 3: 24, 4: 20}
        font_size = sizes.get(block.level, 20)
        height = Inches(0.6)

        txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_text_style(
            run, font_name=fonts["title"], font_size=font_size, bold=True, color=colors["primary"]
        )
        return int(top) + int(height) + Inches(0.1)

    def _render_list(self, slide, block, left, top, width, colors, fonts) -> int:
        """Render bullet or numbered list."""
        items = block.items
        if not items:
            return top

        is_numbered = block.type == ContentType.NUMBERED_LIST
        line_height = Inches(0.35)
        total_height = line_height * len(items) + Inches(0.1)

        txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(total_height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.space_after = Pt(4)
            p.space_before = Pt(2)

            prefix = f"{idx + 1}. " if is_numbered else "\u2022 "
            indent_emu = Inches(0.3 * block.level)
            p.level = min(block.level - 1, 8)

            run = p.add_run()
            run.text = f"{prefix}{item}"
            _set_text_style(run, font_name=fonts["body"], font_size=15, color=colors["text"])

        return int(top) + int(total_height) + Inches(0.15)

    def _render_table(self, slide, block, left, top, width, colors, fonts) -> int:
        """Render a table block."""
        rows_data = block.rows
        if not rows_data:
            return top

        num_rows = len(rows_data)
        num_cols = max(len(r) for r in rows_data) if rows_data else 1
        row_height = Inches(0.4)
        table_height = row_height * num_rows

        # Cap table height
        max_height = Inches(4.5)
        if table_height > max_height:
            table_height = max_height

        tbl_shape = slide.shapes.add_table(
            num_rows, num_cols, int(left), int(top), int(width), int(table_height)
        )
        table = tbl_shape.table

        for row_idx, row_data in enumerate(rows_data):
            for col_idx in range(num_cols):
                cell = table.cell(row_idx, col_idx)
                cell_text = row_data[col_idx] if col_idx < len(row_data) else ""
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(cell_text)

                is_header = row_idx == 0
                if is_header:
                    _set_text_style(
                        run, font_name=fonts["body"], font_size=13, bold=True, color="FFFFFF"
                    )
                    # Header background
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _hex_to_rgb(colors["table_header"])
                else:
                    _set_text_style(
                        run, font_name=fonts["body"], font_size=12, color=colors["text"]
                    )
                    # Alternating row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _hex_to_rgb(colors["table_alt_row"])

                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT

        return int(top) + int(table_height) + Inches(0.25)

    def _render_code(self, slide, block, left, top, width, colors, fonts) -> int:
        """Render a code block with monospace font and background."""
        text = block.content
        if not text:
            return top

        est_lines = max(1, text.count("\n") + 1)
        height = Inches(0.3 * est_lines + 0.3)

        # Background rectangle
        bg_shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb("F5F5F5")
        bg_shape.line.color.rgb = _hex_to_rgb("CCCCCC")
        bg_shape.line.width = Pt(0.5)

        # Code text
        txBox = slide.shapes.add_textbox(
            int(left) + Inches(0.15),
            int(top) + Inches(0.1),
            int(width) - Inches(0.3),
            int(height) - Inches(0.2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_text_style(run, font_name="Consolas", font_size=11, color="333333")

        return int(top) + int(height) + Inches(0.15)

    # ------------------------------------------------------------------ #
    #  Helpers
    # ------------------------------------------------------------------ #

    @staticmethod
    def _set_slide_background(slide, hex_color: str):
        """Set a solid background color for a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)
