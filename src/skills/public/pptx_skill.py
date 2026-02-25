"""PPTX generation skill -- converts various input formats into PowerPoint files."""

from __future__ import annotations

import json
import traceback

from src.generators.base import (
    ContentBlock,
    ContentType,
    GenerationMetadata,
    SlideData,
    StructuredContent,
)
from src.skills.base import BaseSkill
from src.skills.models import QACheck, QAReport, SkillMetadata, SkillResult
from src.utils.logging import get_logger

logger = get_logger(__name__)

# ---------------------------------------------------------------------------
# Content format detection
# ---------------------------------------------------------------------------

def _detect_format(content: str) -> str:
    """Heuristically detect the format of *content*."""
    stripped = content.strip()

    # JSON
    if stripped.startswith("{") or stripped.startswith("["):
        try:
            json.loads(stripped)
            return "json"
        except (json.JSONDecodeError, ValueError):
            pass

    # CSV -- first line has commas and looks tabular
    first_line = stripped.split("\n", 1)[0]
    if "," in first_line and "<" not in first_line:
        comma_count = first_line.count(",")
        if comma_count >= 2:
            return "csv"

    # HTML
    if "<html" in stripped.lower() or stripped.startswith("<!") or "</" in stripped[:200]:
        return "html"

    # Markdown (headings, lists, bold/italic)
    if any(stripped.startswith(c) for c in ("# ", "## ", "### ")):
        return "markdown"
    if "\n# " in stripped or "\n## " in stripped:
        return "markdown"
    if "\n- " in stripped or "\n* " in stripped:
        return "markdown"

    return "text"


# ---------------------------------------------------------------------------
# Lightweight parsers (avoid hard dependency on src.parsers.*)
# ---------------------------------------------------------------------------

def _parse_markdown_to_slides(content: str, title: str) -> StructuredContent:
    """Turn Markdown text into slides split on ``##`` headings."""
    slides: list[SlideData] = []
    current_title = title
    current_blocks: list[ContentBlock] = []
    doc_title = title

    for line in content.split("\n"):
        stripped = line.strip()

        # H1 -> document title (first occurrence only)
        if stripped.startswith("# ") and not stripped.startswith("## "):
            heading = stripped.lstrip("# ").strip()
            if not doc_title:
                doc_title = heading
            continue

        # H2 -> new slide boundary
        if stripped.startswith("## "):
            if current_title or current_blocks:
                slides.append(
                    SlideData(title=current_title, blocks=current_blocks, layout="content")
                )
            current_title = stripped.lstrip("# ").strip()
            current_blocks = []
            continue

        # H3+ -> heading block inside current slide
        if stripped.startswith("### "):
            level = len(stripped.split()[0])  # number of '#'
            current_blocks.append(
                ContentBlock(type=ContentType.HEADING, content=stripped.lstrip("# ").strip(), level=level)
            )
            continue

        # Bullet list item
        if stripped.startswith("- ") or stripped.startswith("* "):
            item_text = stripped[2:].strip()
            # Merge into last block if it's a bullet list, else create one.
            if current_blocks and current_blocks[-1].type == ContentType.BULLET_LIST:
                current_blocks[-1].items.append(item_text)
            else:
                current_blocks.append(
                    ContentBlock(type=ContentType.BULLET_LIST, items=[item_text])
                )
            continue

        # Regular text
        if stripped:
            current_blocks.append(ContentBlock(type=ContentType.TEXT, content=stripped))

    # Flush the last slide.
    if current_title or current_blocks:
        slides.append(SlideData(title=current_title, blocks=current_blocks, layout="content"))

    # Ensure we have at least a title slide.
    if slides and doc_title:
        slides.insert(0, SlideData(title=doc_title, layout="title"))

    return StructuredContent(title=doc_title, slides=slides)


def _parse_text_to_slides(content: str, title: str) -> StructuredContent:
    """Split plain text into slides, one per paragraph."""
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    slides: list[SlideData] = []

    if title:
        slides.append(SlideData(title=title, layout="title"))

    for i, para in enumerate(paragraphs, start=1):
        slides.append(
            SlideData(
                title=f"Slide {i}",
                blocks=[ContentBlock(type=ContentType.TEXT, content=para)],
                layout="content",
            )
        )

    return StructuredContent(title=title or "Untitled", slides=slides)


def _parse_json_to_slides(content: str, title: str) -> StructuredContent:
    """Parse a JSON string into :class:`StructuredContent`.

    Accepts either a raw ``StructuredContent``-shaped dict **or** a simple
    list of ``{"title": ..., "body": ...}`` objects.
    """
    data = json.loads(content)

    # If the JSON already matches StructuredContent shape, parse directly.
    if isinstance(data, dict) and ("slides" in data or "sections" in data):
        sc = StructuredContent.model_validate(data)
        if title:
            sc.title = title
        return sc

    # List of simple slide objects.
    if isinstance(data, list):
        slides = []
        for item in data:
            slide_title = item.get("title", "")
            body = item.get("body", item.get("content", ""))
            blocks = [ContentBlock(type=ContentType.TEXT, content=body)] if body else []
            slides.append(SlideData(title=slide_title, blocks=blocks, layout="content"))
        return StructuredContent(title=title or "Untitled", slides=slides)

    # Fall back to text.
    return _parse_text_to_slides(json.dumps(data, indent=2), title)


def _parse_csv_to_slides(content: str, title: str) -> StructuredContent:
    """Parse CSV content into a single-slide table."""
    rows: list[list[str]] = []
    for line in content.strip().split("\n"):
        if line.strip():
            rows.append([cell.strip() for cell in line.split(",")])

    table_block = ContentBlock(type=ContentType.TABLE, rows=rows)
    slide = SlideData(title=title or "Data", blocks=[table_block], layout="content")
    title_slide = SlideData(title=title or "Data", layout="title")
    return StructuredContent(title=title or "Data", slides=[title_slide, slide])


def _parse_html_to_slides(content: str, title: str) -> StructuredContent:
    """Very lightweight HTML -> slides parser (strips tags)."""
    import re

    text = re.sub(r"<[^>]+>", " ", content)
    text = re.sub(r"\s+", " ", text).strip()
    return _parse_text_to_slides(text, title)


# ---------------------------------------------------------------------------
# The skill
# ---------------------------------------------------------------------------

class PptxSkill(BaseSkill):
    """Skill that generates PPTX (PowerPoint) files."""

    @property
    def metadata(self) -> SkillMetadata:
        return SkillMetadata(
            name="pptx",
            description="Generate PowerPoint presentations from various input formats",
            input_formats=["html", "markdown", "json", "text", "csv"],
            output_format="pptx",
            version="1.0.0",
            tags=["presentation", "office", "pptx"],
        )

    async def validate_params(self, params: dict) -> bool:
        if "content" not in params:
            return False
        if not isinstance(params["content"], str):
            return False
        if not params["content"].strip():
            return False
        return True

    async def execute(self, params: dict, context: dict) -> SkillResult:
        try:
            content_str: str = params["content"]
            content_format: str = params.get("content_format", "") or _detect_format(content_str)
            title: str = params.get("title", "")
            template: str = params.get("template", "default")
            style: dict = params.get("style", {})

            # Parse to StructuredContent based on format.
            parser_map = {
                "markdown": _parse_markdown_to_slides,
                "text": _parse_text_to_slides,
                "json": _parse_json_to_slides,
                "csv": _parse_csv_to_slides,
                "html": _parse_html_to_slides,
            }
            parser = parser_map.get(content_format, _parse_text_to_slides)
            structured = parser(content_str, title)

            gen_meta = GenerationMetadata(template=template, style=style)

            # Try to use the real PptxGenerator; fall back to a minimal one.
            try:
                from src.generators.pptx_generator import PptxGenerator
                generator = PptxGenerator()
            except ImportError:
                logger.warning("pptx_generator_not_found", detail="Using fallback generator")
                generator = _FallbackPptxGenerator()

            output_bytes = await generator.generate(structured, gen_meta)

            return SkillResult(
                success=True,
                output_format="pptx",
                content_bytes=output_bytes,
                metadata={
                    "title": structured.title,
                    "slide_count": len(structured.slides),
                    "content_format": content_format,
                    "template": template,
                },
            )

        except Exception as exc:
            logger.exception("pptx_skill_error")
            return SkillResult(
                success=False,
                output_format="pptx",
                error=f"{type(exc).__name__}: {exc}",
                metadata={"traceback": traceback.format_exc()},
            )

    async def qa_check(self, result: SkillResult) -> QAReport:
        report = await super().qa_check(result)

        # PPTX-specific: check the magic bytes (PK zip header).
        if result.content_bytes:
            is_zip = result.content_bytes[:4] == b"PK\x03\x04"
            report.checks.append(
                QACheck(
                    name="valid_pptx_header",
                    passed=is_zip,
                    detail="File starts with PK zip header" if is_zip else "Invalid PPTX header",
                )
            )
            report.passed = report.passed and is_zip

        return report


# ---------------------------------------------------------------------------
# Fallback generator used when the real one is not yet available.
# ---------------------------------------------------------------------------

class _FallbackPptxGenerator:
    """Minimal PPTX generator using python-pptx directly."""

    async def generate(self, content: StructuredContent, metadata: GenerationMetadata) -> bytes:
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in content.slides:
            if slide_data.layout == "title":
                layout = prs.slide_layouts[0]  # Title Slide
            else:
                layout = prs.slide_layouts[1]  # Title and Content

            slide = prs.slides.add_slide(layout)

            # Title
            if slide.shapes.title and slide_data.title:
                slide.shapes.title.text = slide_data.title

            # Body placeholder
            body_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_placeholder = shape
                    break

            if body_placeholder is not None and slide_data.blocks:
                tf = body_placeholder.text_frame
                tf.clear()
                first = True
                for block in slide_data.blocks:
                    if block.type == ContentType.TEXT or block.type == ContentType.HEADING:
                        if first:
                            tf.text = block.content
                            first = False
                        else:
                            p = tf.add_paragraph()
                            p.text = block.content
                    elif block.type in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
                        for item in block.items:
                            if first:
                                tf.text = item
                                first = False
                            else:
                                p = tf.add_paragraph()
                                p.text = item
                                p.level = 1
                    elif block.type == ContentType.TABLE and block.rows:
                        # Tables don't fit neatly into a placeholder; add text representation.
                        for row in block.rows:
                            if first:
                                tf.text = " | ".join(row)
                                first = False
                            else:
                                p = tf.add_paragraph()
                                p.text = " | ".join(row)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
