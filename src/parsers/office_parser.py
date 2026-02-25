"""Office file parsers: extract StructuredContent from PPTX, DOCX, and XLSX.

Uses:
- python-pptx for PowerPoint files
- python-docx for Word files
- openpyxl for Excel files

Each format is handled by a dedicated method on :class:`OfficeParser`.
"""

from __future__ import annotations

import os
from typing import Optional

from src.generators.base import (
    ContentBlock,
    ContentType,
    SectionData,
    SlideData,
    StructuredContent,
)


class OfficeParser:
    """Parse existing Office files into :class:`StructuredContent`.

    Usage::

        parser = OfficeParser()
        content = parser.parse_file("presentation.pptx")
        content = parser.parse_file("document.docx")
        content = parser.parse_file("spreadsheet.xlsx")
    """

    # Supported extensions mapped to handler names
    _HANDLERS = {
        ".pptx": "_parse_pptx",
        ".docx": "_parse_docx",
        ".xlsx": "_parse_xlsx",
    }

    def parse_file(self, path: str) -> StructuredContent:
        """Detect the file type from its extension and dispatch to the
        appropriate parser method.

        Raises :class:`ValueError` for unsupported extensions.
        """
        ext = os.path.splitext(path)[1].lower()
        handler_name = self._HANDLERS.get(ext)
        if handler_name is None:
            raise ValueError(
                f"Unsupported file extension '{ext}'. "
                f"Supported: {', '.join(sorted(self._HANDLERS))}"
            )
        handler = getattr(self, handler_name)
        return handler(path)

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _parse_pptx(self, path: str) -> StructuredContent:
        """Parse a PowerPoint .pptx file using python-pptx."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        prs = Presentation(path)
        content = StructuredContent()

        # Use the file name (without extension) as a fallback title
        fallback_title = os.path.splitext(os.path.basename(path))[0]

        for slide_idx, slide in enumerate(prs.slides):
            slide_data = SlideData()
            blocks: list[ContentBlock] = []

            for shape in slide.shapes:
                # --- Tables ---
                if shape.has_table:
                    table = shape.table
                    rows: list[list[str]] = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        blocks.append(ContentBlock(type=ContentType.TABLE, rows=rows))
                    continue

                # --- Images ---
                if shape.shape_type is not None:
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image = shape.image
                            blocks.append(
                                ContentBlock(
                                    type=ContentType.IMAGE,
                                    content=shape.name or "",
                                    metadata={
                                        "width": shape.width,
                                        "height": shape.height,
                                        "content_type": image.content_type if image else "",
                                    },
                                )
                            )
                            continue
                    except (ImportError, AttributeError):
                        pass

                # --- Text frames ---
                if shape.has_text_frame:
                    tf = shape.text_frame
                    for paragraph in tf.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue

                        # Determine if this is a title placeholder
                        is_title = False
                        try:
                            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                                ph_idx = shape.placeholder_format.idx
                                # idx 0 = title, idx 1 = subtitle/body
                                if ph_idx == 0:
                                    is_title = True
                                elif ph_idx == 1 and not slide_data.title:
                                    # Could be subtitle
                                    if not slide_data.subtitle:
                                        slide_data.subtitle = text
                                    continue
                        except Exception:
                            pass

                        if is_title and not slide_data.title:
                            slide_data.title = text
                        elif paragraph.level and paragraph.level > 0:
                            # Indented text -> bullet item
                            blocks.append(
                                ContentBlock(
                                    type=ContentType.BULLET_LIST,
                                    items=[text],
                                )
                            )
                        else:
                            blocks.append(
                                ContentBlock(type=ContentType.TEXT, content=text)
                            )

            slide_data.blocks = blocks
            if not slide_data.title and blocks:
                # Use first text block as title if none was found
                for b in blocks:
                    if b.type == ContentType.TEXT and b.content:
                        slide_data.title = b.content
                        break

            layout = "title" if (not blocks and slide_data.title) else "content"
            slide_data.layout = layout

            content.slides.append(slide_data)

            # Mirror into sections
            content.sections.append(
                SectionData(title=slide_data.title, blocks=list(blocks))
            )

        # Set overall title from first slide
        if content.slides:
            content.title = content.slides[0].title or fallback_title
        else:
            content.title = fallback_title

        return content

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _parse_docx(self, path: str) -> StructuredContent:
        """Parse a Word .docx file using python-docx."""
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document(path)
        content = StructuredContent()

        fallback_title = os.path.splitext(os.path.basename(path))[0]

        # Track current section being built
        current_blocks: list[ContentBlock] = []
        section_title = ""
        first_heading_seen = False

        def _commit_section() -> None:
            nonlocal current_blocks, section_title
            if section_title or current_blocks:
                content.sections.append(
                    SectionData(title=section_title, blocks=list(current_blocks))
                )
                content.slides.append(
                    SlideData(
                        title=section_title,
                        blocks=list(current_blocks),
                        layout="content" if current_blocks else "title",
                    )
                )
            current_blocks = []
            section_title = ""

        # --- Paragraphs ---
        for para in doc.paragraphs:
            style_name = (para.style.name or "").lower() if para.style else ""
            text = para.text.strip()

            if not text:
                continue

            # Heading styles
            if style_name.startswith("heading"):
                # Extract level from style name (e.g. "Heading 1" -> 1)
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 1

                if level == 1:
                    if first_heading_seen:
                        _commit_section()
                    else:
                        first_heading_seen = True
                        content.title = text
                    section_title = text
                else:
                    current_blocks.append(
                        ContentBlock(type=ContentType.HEADING, content=text, level=level)
                    )

            elif style_name.startswith("list bullet") or style_name == "list bullet":
                # Bullet list item
                current_blocks.append(
                    ContentBlock(type=ContentType.BULLET_LIST, items=[text])
                )

            elif style_name.startswith("list number") or style_name == "list number":
                # Numbered list item
                current_blocks.append(
                    ContentBlock(type=ContentType.NUMBERED_LIST, items=[text])
                )

            elif style_name == "title":
                if not content.title:
                    content.title = text
                if not section_title:
                    section_title = text

            elif style_name == "subtitle":
                content.subtitle = text

            else:
                current_blocks.append(
                    ContentBlock(type=ContentType.TEXT, content=text)
                )

        # --- Tables ---
        for table in doc.tables:
            rows: list[list[str]] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                current_blocks.append(ContentBlock(type=ContentType.TABLE, rows=rows))

        # Commit final section
        _commit_section()

        if not content.title:
            if content.sections:
                content.title = content.sections[0].title or fallback_title
            else:
                content.title = fallback_title

        return content

    # ------------------------------------------------------------------
    # XLSX
    # ------------------------------------------------------------------

    def _parse_xlsx(self, path: str) -> StructuredContent:
        """Parse an Excel .xlsx file using openpyxl."""
        from openpyxl import load_workbook

        wb = load_workbook(path, read_only=True, data_only=True)
        content = StructuredContent()

        fallback_title = os.path.splitext(os.path.basename(path))[0]
        content.title = fallback_title

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            all_rows: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                all_rows.append([str(cell) if cell is not None else "" for cell in row])

            if not all_rows:
                # Empty sheet -- still record it
                content.sheets.append({"name": sheet_name, "headers": [], "rows": []})
                continue

            headers = all_rows[0]
            data_rows = all_rows[1:]

            # Populate sheets for XLSX round-trip
            content.sheets.append(
                {
                    "name": sheet_name,
                    "headers": headers,
                    "rows": data_rows,
                }
            )

            # Create a table block for slide/section output
            table_block = ContentBlock(
                type=ContentType.TABLE,
                rows=[headers] + data_rows,
            )

            content.slides.append(
                SlideData(
                    title=sheet_name,
                    blocks=[table_block],
                    layout="content",
                )
            )
            content.sections.append(
                SectionData(title=sheet_name, blocks=[table_block])
            )

        wb.close()
        return content
